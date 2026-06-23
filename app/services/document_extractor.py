import io
import os
import logging
import tempfile
from urllib.parse import urlparse

import httpx

logger = logging.getLogger(__name__)

# Supported document extensions and their MIME types
DOCUMENT_EXTENSIONS = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".doc": "application/msword",
    ".xls": "application/vnd.ms-excel",
    ".ppt": "application/vnd.ms-powerpoint",
    ".txt": "text/plain",
    ".csv": "text/csv",
}

DOCUMENT_MIME_TYPES = {v: k for k, v in DOCUMENT_EXTENSIONS.items()}


def is_document_url(url: str) -> bool:
    """Check if a URL likely points to a downloadable document."""
    parsed = urlparse(url)
    path = parsed.path.lower()
    return any(path.endswith(ext) for ext in DOCUMENT_EXTENSIONS)


def detect_document_type(url: str, content_type: str | None = None) -> str | None:
    """Detect the document type from URL extension or Content-Type header."""
    parsed = urlparse(url)
    path = parsed.path.lower()

    for ext in DOCUMENT_EXTENSIONS:
        if path.endswith(ext):
            return ext

    if content_type:
        ct = content_type.lower().split(";")[0].strip()
        if ct in DOCUMENT_MIME_TYPES:
            return DOCUMENT_MIME_TYPES[ct]

    return None


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF file."""
    try:
        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(file_bytes))
        text_parts = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text.strip())
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error("PDF extraction failed: %s", e)
        return ""


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from a DOCX file."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())

        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_data:
                    text_parts.append(" | ".join(row_data))

        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error("DOCX extraction failed: %s", e)
        return ""


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    """Extract text from an XLSX file."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) for cell in row if cell is not None]
                if row_data:
                    text_parts.append(" | ".join(row_data))
        wb.close()
        return "\n".join(text_parts)
    except Exception as e:
        logger.error("XLSX extraction failed: %s", e)
        return ""


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from a PPTX file."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            slide_texts.append(para_text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_data:
                            slide_texts.append(" | ".join(row_data))
            if slide_texts:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error("PPTX extraction failed: %s", e)
        return ""


def extract_text_from_txt(file_bytes: bytes) -> str:
    """Extract text from plain text or CSV files."""
    try:
        for encoding in ["utf-8", "latin-1", "cp1252"]:
            try:
                return file_bytes.decode(encoding)
            except UnicodeDecodeError:
                continue
        return file_bytes.decode("utf-8", errors="replace")
    except Exception as e:
        logger.error("Text extraction failed: %s", e)
        return ""


EXTRACTORS = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".doc": extract_text_from_docx,  # Best-effort, may not work for old .doc
    ".xlsx": extract_text_from_xlsx,
    ".xls": extract_text_from_xlsx,
    ".pptx": extract_text_from_pptx,
    ".ppt": extract_text_from_pptx,
    ".txt": extract_text_from_txt,
    ".csv": extract_text_from_txt,
}


def detect_extension_from_filename(filename: str) -> str | None:
    """Detect supported document extension from a filename."""
    lower = filename.lower().strip()
    for ext in DOCUMENT_EXTENSIONS:
        if lower.endswith(ext):
            return ext
    return None


def extract_text_from_bytes(file_bytes: bytes, extension: str) -> str:
    """Extract text from uploaded file bytes. Raises ValueError if unsupported."""
    ext = extension.lower()
    if not ext.startswith("."):
        ext = f".{ext}"

    extractor = EXTRACTORS.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported document type: {ext}")

    text = extractor(file_bytes)
    if not text or not text.strip():
        raise ValueError(f"No text could be extracted from {ext} file.")

    return text


async def download_and_extract_document(url: str, timeout: int = 60) -> tuple[str | None, str | None, str]:
    """
    Download a document from a URL and extract its text.

    Returns:
        (title, text, doc_type) where doc_type is the file extension.
    """
    try:
        async with httpx.AsyncClient(
            timeout=timeout,
            follow_redirects=True,
            headers={
                "User-Agent": (
                    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                    "AppleWebKit/537.36 (KHTML, like Gecko) "
                    "Chrome/122.0.0.0 Safari/537.36"
                )
            },
        ) as client:
            response = await client.get(url)
            response.raise_for_status()

            content_type = response.headers.get("content-type", "")
            doc_type = detect_document_type(url, content_type)

            if not doc_type:
                logger.warning("Could not determine document type for %s", url)
                return None, None, ""

            file_bytes = response.content

            extractor = EXTRACTORS.get(doc_type)
            if not extractor:
                logger.warning("No extractor available for type: %s", doc_type)
                return None, None, doc_type

            text = extractor(file_bytes)

            # Generate a title from the filename
            parsed_path = urlparse(url).path
            filename = os.path.basename(parsed_path)
            title = filename if filename else f"Document ({doc_type})"

            logger.info(
                "Document extraction (%s): %d chars from %s",
                doc_type, len(text), url
            )

            return title, text, doc_type

    except Exception as e:
        logger.error("Document download/extraction failed for %s: %s", url, e)
        return None, None, ""
